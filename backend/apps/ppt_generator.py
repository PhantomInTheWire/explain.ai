import json
import asyncio
from pathlib import Path

from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Pt

from apps.promptops import process_user_question, set_session_json_slide
from core.vectorstore import vectorstore_manager
from core.storage import storage_manager
from core.logging import get_logger

log = get_logger(__name__)

THEME_CONFIG = {
    "Theme1": (0, 1),
    "Theme2": (11, 14),
    "Theme3": (0, 10),
    "Theme4": (0, 10),
}
THEME_BASE_PATH = Path(__file__).parent.parent / "theme_pptx"


async def generate_presentation(
    session_id: str, selected_theme: str = "Theme1"
) -> Path:
    if not await vectorstore_manager.collection_exists(session_id):
        raise ValueError("Please upload a PDF file first.")

    prompt = """Write a powerpoint presentation about this document in JSON.
The JSON should contain 'slides' and each slide should contain 'title' and 'content'.
Use bullet points to make text engaging. Divide content into multiple slides if needed.
Output only valid JSON with no explanations."""

    try:
        response_str = await process_user_question(session_id, prompt)
        json_string = response_str.strip()

        for prefix in ["```json", "```"]:
            if json_string.startswith(prefix):
                json_string = json_string[len(prefix) :]
        if json_string.endswith("```"):
            json_string = json_string[:-3]

        json_pptx = json.loads(json_string.strip())
        await set_session_json_slide(session_id, json_pptx)
        return await process_presentation(session_id, json_pptx, selected_theme)
    except json.JSONDecodeError as e:
        log.error("json parsing error", error=str(e))
        raise ValueError("Failed to parse presentation structure")
    except Exception as e:
        log.error("presentation generation error", error=str(e))
        raise


def validate_json_data_structure(json_data) -> None:
    if not isinstance(json_data, (dict, list)):
        raise ValueError("Invalid JSON data format")
    slides = json_data.get("slides") if isinstance(json_data, dict) else json_data
    if slides is None:
        raise ValueError("Missing 'slides' field")
    for slide in slides:
        if (
            not isinstance(slide, dict)
            or "title" not in slide
            or "content" not in slide
        ):
            raise ValueError("Invalid slide structure")


def apply_bold_format(word: str, r, bolded_mode: bool) -> bool:
    if word.startswith("**") and word.endswith("**"):
        r.text = word[2:-2]
        r.font.bold = True
    elif word.startswith("**"):
        r.text = word[2:] + " "
        r.font.bold = True
        bolded_mode = True
    elif word.endswith("**"):
        r.text = word[:-2] + " "
        r.font.bold = True
        bolded_mode = False
    elif bolded_mode:
        r.text = word + " "
        r.font.bold = True
    else:
        r.text = word + " "
    return bolded_mode


def create_slide(
    prs: Presentation, title: str, content: str, selected_theme: str
) -> None:
    layout, text_placeholder = THEME_CONFIG.get(selected_theme, (11, 14))
    slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    content_shape = slide.placeholders[text_placeholder]
    content_text_frame = content_shape.text_frame
    content_text_frame.word_wrap = True

    for line in content.replace("\\n", "\n").split("\n"):
        p = content_text_frame.add_paragraph()
        if line.startswith("##"):
            p.text = line[2:].strip()
            p.font.size = Pt(20)
            p.font.bold = True
        elif line.startswith("* "):
            words = line[1:].strip().split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)
            for run in p.runs:
                run.font.size = Pt(12)
        else:
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)
                r.font.size = Pt(14)


async def process_presentation(
    session_id: str, json_data: dict, selected_theme: str = "Theme1"
) -> Path:
    # python-pptx is synchronous
    def _process():
        validate_json_data_structure(json_data)

        theme_path = THEME_BASE_PATH / f"{selected_theme}.pptx"
        if not theme_path.exists():
            theme_path = THEME_BASE_PATH / "Theme1.pptx"

        presentation = Presentation(str(theme_path))
        slides = (
            json_data.get("slides", []) if isinstance(json_data, dict) else json_data
        )

        for slide_data in slides:
            if isinstance(slide_data, dict):
                create_slide(
                    presentation,
                    slide_data.get("title", ""),
                    slide_data.get("content", "").replace("-", ""),
                    selected_theme,
                )

        output_path = storage_manager.get_output_path(session_id, "presentation.pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))

        log.info("presentation generated", output_path=str(output_path))
        return output_path

    return await asyncio.to_thread(_process)
